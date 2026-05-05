"""
_build_pptx.py — generates greek_banking_deck.pptx (12 slides, board-ready)
Run: python deliverables/_build_pptx.py
"""
import sys, io, sqlite3
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.gridspec import GridSpec
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.stdout.reconfigure(encoding="utf-8")

ROOT = Path(__file__).resolve().parent.parent
DB   = ROOT / "02_Banking_Sector_Dashboard" / "data" / "greek_banking_final.db"
OUT  = Path(__file__).parent / "greek_banking_deck.pptx"

# ── Data ──────────────────────────────────────────────────────────────────────
con  = sqlite3.connect(DB)
kpis = pd.read_sql("SELECT * FROM kpis ORDER BY bank, year", con)
con.close()

BANKS = ["Eurobank", "Alpha Bank", "Piraeus Bank", "NBG"]
YEARS = [2022, 2023, 2024]

BCOLORS = {
    "Eurobank":     (0,   103, 177),
    "Alpha Bank":   (226,  35,  26),
    "Piraeus Bank": (247, 166,   0),
    "NBG":          (  0,  48, 135),
}
BCOLORS_HEX = {
    "Eurobank":     RGBColor(0,   103, 177),
    "Alpha Bank":   RGBColor(226,  35,  26),
    "Piraeus Bank": RGBColor(247, 166,   0),
    "NBG":          RGBColor(  0,  48, 135),
}

NAVY  = RGBColor(10,  22,  40)
WHITE = RGBColor(255,255,255)
LGREY = RGBColor(148,163,184)
DGREY = RGBColor(30, 41, 59)
GREEN = RGBColor(16, 163, 127)


# ── Helpers ───────────────────────────────────────────────────────────────────
def rgb(r, g, b):
    return RGBColor(r, g, b)

def hex2rgb(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def fig_to_stream(fig, dpi=150):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf

MPL_THEME = {
    "figure.facecolor": "#0a1628",
    "axes.facecolor":   "#0f1e30",
    "axes.edgecolor":   "#1e293b",
    "axes.labelcolor":  "#94a3b8",
    "axes.titlecolor":  "#f1f5f9",
    "xtick.color":      "#64748b",
    "ytick.color":      "#64748b",
    "text.color":       "#94a3b8",
    "grid.color":       "#1e293b",
    "grid.alpha":       1.0,
    "lines.linewidth":  2.0,
    "font.family":      "sans-serif",
    "font.size":        9,
}

def mpl_context():
    return plt.rc_context(MPL_THEME)

def bank_color_mpl(bank):
    r, g, b = BCOLORS[bank]
    return (r/255, g/255, b/255)


# ── Presentation setup ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]   # completely blank

# ── Slide factory helpers ──────────────────────────────────────────────────────
def add_slide():
    sl = prs.slides.add_slide(blank_layout)
    # Dark background
    bg = sl.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    return sl

def add_rect(sl, x, y, w, h, color):
    shape = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(sl, text, x, y, w, h, size=12, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb

def add_image(sl, stream, x, y, w, h=None):
    if h:
        sl.shapes.add_picture(stream, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        sl.shapes.add_picture(stream, Inches(x), Inches(y), Inches(w))

def slide_header(sl, title, subtitle=None, bar_color=DGREY):
    add_rect(sl, 0, 0, 13.33, 1.1, bar_color)
    add_text(sl, title, 0.35, 0.12, 11, 0.7, size=22, bold=True, color=WHITE)
    if subtitle:
        add_text(sl, subtitle, 0.35, 0.75, 11, 0.3, size=10, color=LGREY)

def footer(sl, page_num):
    add_text(sl, f"Greek Banking Sector Analysis 2022–2024  |  Spyros Papastergiou  |  {page_num}",
             0.3, 7.2, 12.5, 0.25, size=8, color=LGREY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
sl1 = add_slide()
add_rect(sl1, 0, 0, 13.33, 7.5, NAVY)
add_rect(sl1, 0, 0, 0.25, 7.5, BCOLORS_HEX["Eurobank"])   # left accent bar

add_text(sl1, "GREEK BANKING SECTOR ANALYSIS",
         0.6, 1.8, 12, 1.0, size=32, bold=True, color=WHITE)
add_text(sl1, "2022 – 2024",
         0.6, 2.85, 6, 0.6, size=24, bold=False, color=LGREY)
add_text(sl1,
         "Eurobank  ·  Alpha Bank  ·  Piraeus Bank  ·  NBG",
         0.6, 3.6, 10, 0.5, size=14, color=LGREY)

add_rect(sl1, 0.6, 4.4, 3.5, 0.04, BCOLORS_HEX["Eurobank"])

add_text(sl1, "Spyros Papastergiou", 0.6, 4.6, 6, 0.4, size=12, bold=True, color=WHITE)
add_text(sl1, "spyrossyo96@gmail.com  |  github.com/papastergiousp-maker/greek-banking-sector-analysis",
         0.6, 5.05, 11, 0.35, size=9, color=LGREY)
add_text(sl1, "Data sourced exclusively from audited annual reports · Not investment advice",
         0.6, 6.8, 11, 0.35, size=8, color=rgb(71,85,105))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ════════════════════════════════════════════════════════════════════════════
sl2 = add_slide()
slide_header(sl2, "Executive Summary",
             "Key findings from 12 annual reports  |  FY2022–2024")
footer(sl2, "2 / 12")

findings = [
    (BCOLORS_HEX["Eurobank"],
     "+55% Sector NII Growth",
     "Net interest income rose from €5.5bn (2022) to €8.6bn (2024) driven "
     "almost entirely by ECB rate hikes. NBG leads with NIM of 3.14% — "
     "highest in the peer group."),
    (BCOLORS_HEX["NBG"],
     "NPE Cleanup Complete",
     "All four banks reduced NPE ratios below 4% by 2024 (from 5–8% in 2022). "
     "NBG and Piraeus tied at 2.6% — a structural improvement that took a decade "
     "and €tens of billions of portfolio disposals."),
    (BCOLORS_HEX["Piraeus Bank"],
     "Rate Cycle Now a Headwind",
     "ECB cuts beginning in 2024 will reduce sector NII by ~€170–220m per 25bp. "
     "Under EBA-style stress (+200bp CoR, −15% loans, −50bp NIM), Piraeus Bank's "
     "CET1 falls to ~9.9% — breaching the 10.5% SREP floor."),
]

for i, (col, title, body) in enumerate(findings):
    x = 0.35 + i * 4.3
    add_rect(sl2, x, 1.3, 4.1, 0.06, col)
    add_text(sl2, title, x, 1.45, 4.0, 0.45, size=13, bold=True, color=WHITE)
    add_text(sl2, body,  x, 1.95, 4.0, 1.8,  size=9.5, color=LGREY)

# KPI summary table
headers = ["Bank", "ROE 2024", "NIM 2024", "CET1 2024", "NPE 2024", "Net Profit 2024"]
col_w   = [2.2, 1.6, 1.6, 1.6, 1.5, 2.1]
col_x   = [0.35]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

add_rect(sl2, 0.35, 4.1, sum(col_w), 0.35, DGREY)
for j, (hdr, cx) in enumerate(zip(headers, col_x)):
    add_text(sl2, hdr, cx+0.05, 4.12, col_w[j]-0.1, 0.3,
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

k24 = kpis[kpis.year == 2024].set_index("bank")
for i, bank in enumerate(BANKS):
    row_y = 4.5 + i * 0.42
    row_bg = rgb(15,30,50) if i % 2 == 0 else DGREY
    add_rect(sl2, 0.35, row_y, sum(col_w), 0.38, row_bg)
    vals = [
        bank,
        f"{k24.loc[bank,'roe']:.1f}%",
        f"{k24.loc[bank,'nim']:.2f}%",
        f"{k24.loc[bank,'cet1']:.1f}%",
        f"{k24.loc[bank,'npe_ratio']:.1f}%",
        f"€{k24.loc[bank,'net_profit']:,.0f}m",
    ]
    for j, (val, cx) in enumerate(zip(vals, col_x)):
        is_bank = j == 0
        add_text(sl2, val, cx+0.05, row_y+0.04, col_w[j]-0.1, 0.32,
                 size=9.5, bold=is_bank,
                 color=RGBColor(*BCOLORS[bank]) if is_bank else WHITE,
                 align=PP_ALIGN.LEFT if is_bank else PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SECTOR LANDSCAPE
# ════════════════════════════════════════════════════════════════════════════
sl3 = add_slide()
slide_header(sl3, "Sector Landscape — NII & Profitability",
             "Net Interest Income surge driven by ECB rate cycle  |  All values €m")
footer(sl3, "3 / 12")

with mpl_context():
    fig, axes = plt.subplots(1, 2, figsize=(11, 4.5))
    fig.patch.set_facecolor("#0a1628")

    # Left: NII grouped bar
    ax = axes[0]
    x   = np.arange(len(BANKS))
    w   = 0.25
    for j, yr in enumerate(YEARS):
        vals = [kpis[(kpis.bank==b)&(kpis.year==yr)]["nii"].values[0] for b in BANKS]
        bars = ax.bar(x + (j-1)*w, vals, w*0.92,
                      color=[bank_color_mpl(b) for b in BANKS],
                      alpha=0.6+j*0.15, label=str(yr))
        for bar, v in zip(bars, vals):
            if yr == 2024:
                ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+20,
                        f"€{v/1000:.1f}bn", ha="center", va="bottom", fontsize=6.5,
                        color="#f1f5f9")
    ax.set_title("Net Interest Income by Bank (€m)", fontsize=10, pad=8)
    ax.set_xticks(x)
    ax.set_xticklabels(BANKS, fontsize=8)
    ax.legend(fontsize=8, framealpha=0)
    ax.grid(axis="y", alpha=0.3)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"€{v/1000:.1f}bn" if v >= 1000 else f"€{v:,.0f}m"))

    # Right: ROE trend lines
    ax2 = axes[1]
    for bank in BANKS:
        bdf = kpis[kpis.bank == bank].sort_values("year")
        ax2.plot(bdf["year"], bdf["roe"],
                 color=bank_color_mpl(bank), marker="o", markersize=6, label=bank)
        ax2.annotate(f"{bdf[bdf.year==2024]['roe'].values[0]:.1f}%",
                     xy=(2024, bdf[bdf.year==2024]["roe"].values[0]),
                     xytext=(5, 0), textcoords="offset points",
                     fontsize=7.5, color=bank_color_mpl(bank))
    ax2.axhline(11, linestyle="--", color="#475569", linewidth=1, label="CoE ~11%")
    ax2.set_title("Return on Equity (%)", fontsize=10, pad=8)
    ax2.set_xticks(YEARS)
    ax2.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"{v:.0f}%"))
    ax2.legend(fontsize=7.5, framealpha=0)
    ax2.grid(axis="y", alpha=0.3)

    fig.tight_layout(pad=1.5)
    stream3 = fig_to_stream(fig, dpi=150)

add_image(sl3, stream3, 0.35, 1.2, 12.5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDES 4–7 — BANK ONE-PAGERS
# ════════════════════════════════════════════════════════════════════════════
BANK_BULLETS = {
    "Eurobank": [
        "ROE of 16.9% — highest in peer group; only bank trading below justified P/B (~1.5x vs market ~1.3x)",
        "Hellenic Bank consolidation adds Cypriot diversification; 71% of loans outside Greece by 2024",
        "Best earnings quality: lowest C/I in sector (32.6%); 710bp CET1 headroom above SREP floor",
    ],
    "Alpha Bank": [
        "ROE of 8.2% below CoE of ~11% — justified P/B ~0.7x; re-rate requires ROE expansion",
        "NPE cleanup accelerated: ratio fell from 7.8% (2022) to 3.8% (2024) via GAIA II and Project ACAC",
        "NII flat in 2024 (−0.3% YoY) as rate tailwind fades; operational leverage improving (C/I 37.7%)",
    ],
    "Piraeus Bank": [
        "2024 profit surge (+35% YoY) partly driven by ~€165m after-tax impairment benefit; underlying growth ~+14%",
        "Best-in-class efficiency: C/I of 31.8% — lowest in sector; strong operational improvement trajectory",
        "Tightest capital buffer: CET1 13.7%, only 320bp above SREP floor; breaches in EBA adverse stress",
    ],
    "NBG": [
        "Highest NIM (3.14%) and CET1 (19.1%) in sector; 860bp headroom above SREP floor",
        "NPE ratio of 2.6% — tied for sector-best; ROE of 13.7% sustained above CoE despite rate headwinds",
        "Structural rate sensitivity: higher fixed-spread SME/retail mix provides partial ECB rate hedge",
    ],
}

for pg, bank in enumerate(BANKS, 4):
    sl = add_slide()
    bcolor = BCOLORS_HEX[bank]
    slide_header(sl, bank, f"FY2024 Snapshot  |  {BANKS.index(bank)+1} of 4", bar_color=bcolor)
    footer(sl, f"{pg} / 12")

    k24b = kpis[(kpis.bank==bank)&(kpis.year==2024)].iloc[0]
    k23b = kpis[(kpis.bank==bank)&(kpis.year==2023)].iloc[0]
    k22b = kpis[(kpis.bank==bank)&(kpis.year==2022)].iloc[0]

    # KPI cards row
    metrics = [
        ("ROE",   f"{k24b['roe']:.1f}%",   f"{k24b['roe']-k23b['roe']:+.1f}pp"),
        ("NIM",   f"{k24b['nim']:.2f}%",   f"{k24b['nim']-k23b['nim']:+.2f}pp"),
        ("CET1",  f"{k24b['cet1']:.1f}%",  f"{k24b['cet1']-k23b['cet1']:+.1f}pp"),
        ("C/I",   f"{k24b['cost_to_income']:.1f}%", f"{k24b['cost_to_income']-k23b['cost_to_income']:+.1f}pp"),
        ("NPE",   f"{k24b['npe_ratio']:.1f}%", f"{k24b['npe_ratio']-k23b['npe_ratio']:+.1f}pp"),
        ("Net\nProfit", f"€{k24b['net_profit']:,.0f}m",
         f"{(k24b['net_profit']-k23b['net_profit'])/abs(k23b['net_profit'])*100:+.0f}%"),
    ]
    cw = 1.9
    for mi, (lbl, val, delta) in enumerate(metrics):
        cx = 0.35 + mi * (cw + 0.1)
        add_rect(sl, cx, 1.25, cw, 1.0, DGREY)
        add_rect(sl, cx, 1.25, cw, 0.06, bcolor)
        add_text(sl, lbl,   cx+0.05, 1.32, cw-0.1, 0.3, size=9,  color=LGREY, align=PP_ALIGN.CENTER)
        add_text(sl, val,   cx+0.05, 1.58, cw-0.1, 0.45, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        dcolor = GREEN if not delta.startswith("-") else rgb(239,68,68)
        if lbl in ("C/I", "NPE"):
            dcolor = rgb(239,68,68) if not delta.startswith("-") else GREEN
        add_text(sl, delta, cx+0.05, 2.0,  cw-0.1, 0.22, size=9, color=dcolor, align=PP_ALIGN.CENTER)

    # Bullet points
    add_rect(sl, 0.35, 2.45, 6.5, 0.04, bcolor)
    add_text(sl, "Investment Highlights", 0.35, 2.55, 6.5, 0.35,
             size=11, bold=True, color=WHITE)
    for bi, bullet in enumerate(BANK_BULLETS[bank]):
        add_text(sl, f"▸  {bullet}", 0.35, 3.0 + bi * 0.75, 6.5, 0.7,
                 size=9.5, color=LGREY)

    # Trend chart (NII + ROE dual-axis)
    with mpl_context():
        fig, ax1 = plt.subplots(figsize=(5, 3.5))
        fig.patch.set_facecolor("#0a1628")
        bdf = kpis[kpis.bank==bank].sort_values("year")
        bc  = bank_color_mpl(bank)
        ax1.bar(bdf["year"], bdf["nii"], 0.4, color=bc, alpha=0.65, label="NII")
        ax1.set_ylabel("NII (€m)", fontsize=8)
        ax1.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"€{v:,.0f}m"))
        ax2 = ax1.twinx()
        ax2.plot(bdf["year"], bdf["roe"], color="#10b981", marker="o", markersize=6, linewidth=2, label="ROE")
        ax2.axhline(11, linestyle="--", color="#475569", linewidth=1)
        ax2.set_ylabel("ROE (%)", fontsize=8, color="#10b981")
        ax2.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"{v:.0f}%"))
        ax2.tick_params(axis="y", colors="#10b981")
        ax1.set_xticks(YEARS)
        ax1.set_title(f"{bank} — NII & ROE", fontsize=9, pad=6)
        ax1.grid(axis="y", alpha=0.2)
        lines1, labels1 = ax1.get_legend_handles_labels()
        lines2, labels2 = ax2.get_legend_handles_labels()
        ax1.legend(lines1+lines2, labels1+labels2, fontsize=7.5, framealpha=0, loc="upper left")
        fig.tight_layout()
        stream_bank = fig_to_stream(fig, dpi=130)

    add_image(sl, stream_bank, 7.1, 2.4, 5.8)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CROSS-BANK RANKINGS
# ════════════════════════════════════════════════════════════════════════════
sl8 = add_slide()
slide_header(sl8, "Peer Rankings — FY2024",
             "Green = best in class per metric  |  Direction-adjusted (lower C/I, NPE = better)")
footer(sl8, "8 / 12")

with mpl_context():
    fig, axes = plt.subplots(1, 3, figsize=(12, 4.5))
    fig.patch.set_facecolor("#0a1628")

    ranking_metrics = [
        ("roe",           "ROE (%)",           True,  "%"),
        ("nim",           "NIM (%)",            True,  "%"),
        ("cost_to_income","Cost-to-Income (%)",False, "%"),
        ("cet1",          "CET1 (%)",           True,  "%"),
        ("npe_ratio",     "NPE Ratio (%)",      False, "%"),
        ("net_profit",    "Net Profit (€m)",    True,  ""),
    ]

    for ai, (col_key, label, higher_better, sfx) in enumerate(ranking_metrics):
        ax = axes[ai % 3]
        if ai == 3:
            ax = axes[0]

    # Use a 2x3 grid of small charts
    fig2, axes2 = plt.subplots(2, 3, figsize=(12.5, 5))
    fig2.patch.set_facecolor("#0a1628")
    plt.subplots_adjust(hspace=0.55, wspace=0.35)

    for ai, (col_key, label, higher_better, sfx) in enumerate(ranking_metrics):
        ax = axes2[ai//3, ai%3]
        k24_df = kpis[kpis.year==2024].set_index("bank")
        vals   = k24_df[col_key].reindex(BANKS)
        sorted_banks = vals.sort_values(ascending=not higher_better).index
        sorted_vals  = vals[sorted_banks]
        colors_sorted = [bank_color_mpl(b) for b in sorted_banks]
        bars = ax.barh(range(len(sorted_banks)), sorted_vals.values,
                       color=colors_sorted, height=0.6)
        ax.set_yticks(range(len(sorted_banks)))
        ax.set_yticklabels([b.replace(" Bank","").replace("Piraeus","Piraeus") for b in sorted_banks],
                           fontsize=7.5)
        ax.set_title(label, fontsize=8.5, pad=4)
        ax.grid(axis="x", alpha=0.3)
        for bar, v in zip(bars, sorted_vals.values):
            fmt = f"{v:.1f}{sfx}" if sfx == "%" else f"€{v:,.0f}m"
            ax.text(v + (sorted_vals.max()*0.01), bar.get_y()+bar.get_height()/2,
                    fmt, va="center", fontsize=7, color="#f1f5f9")
        ax.set_xlim(0, sorted_vals.max() * 1.2)
        # Highlight best bar
        bars[0].set_edgecolor("#10b981")
        bars[0].set_linewidth(1.5)

    fig2.tight_layout(pad=1.0)
    stream8 = fig_to_stream(fig2, dpi=130)
    plt.close(fig)

add_image(sl8, stream8, 0.3, 1.2, 12.7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — NII FORECAST
# ════════════════════════════════════════════════════════════════════════════
sl9 = add_slide()
slide_header(sl9, "NII Forecast 2025–2026 — ECB Rate Scenarios",
             "Rate sensitivity: ~€195m sector NII per 25bp ECB cut  |  Volume: +2% p.a. assumed")
footer(sl9, "9 / 12")

ECB_HIST = {2022: 0.56, 2023: 3.39, 2024: 3.28}
SCENARIOS = {
    "Dovish":  {2025: 1.75, 2026: 1.50},
    "Base":    {2025: 2.25, 2026: 2.00},
    "Hawkish": {2025: 3.15, 2026: 3.15},
}
SC_COLORS = {"Dovish": "#60a5fa", "Base": "#10b981", "Hawkish": "#f59e0b"}
NII_SENS   = -195  # €m per 25bp sector cut

with mpl_context():
    fig, (ax_rate, ax_nii) = plt.subplots(1, 2, figsize=(12, 4.5))
    fig.patch.set_facecolor("#0a1628")

    # ECB rate path
    all_years = [2022, 2023, 2024, 2025, 2026]
    for sc_name, sc_rates in SCENARIOS.items():
        rates = [ECB_HIST[2022], ECB_HIST[2023], ECB_HIST[2024],
                 sc_rates[2025], sc_rates[2026]]
        lw = 2.5 if sc_name == "Base" else 1.5
        ls = "-" if sc_name == "Base" else "--"
        ax_rate.plot(all_years, rates, color=SC_COLORS[sc_name],
                     linewidth=lw, linestyle=ls, marker="o", markersize=5, label=sc_name)
    ax_rate.axvline(2024.5, linestyle=":", color="#475569", linewidth=1)
    ax_rate.set_title("ECB Deposit Facility Rate Path (%)", fontsize=9.5)
    ax_rate.set_xticks(all_years)
    ax_rate.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"{v:.1f}%"))
    ax_rate.legend(fontsize=8, framealpha=0)
    ax_rate.grid(alpha=0.3)

    # NII forecast per bank — base scenario
    k24_nii    = kpis[kpis.year==2024].set_index("bank")["nii"]
    sector_nii = k24_nii.sum()

    for bank in BANKS:
        hist_vals = [kpis[(kpis.bank==bank)&(kpis.year==y)]["nii"].values[0] for y in YEARS]
        ax_nii.plot(YEARS, hist_vals, color=bank_color_mpl(bank),
                    linewidth=2, marker="o", markersize=5)

        nii_prev = k24_nii[bank]
        fc_vals  = [nii_prev]
        for y in [2025, 2026]:
            prev_r = ECB_HIST[2024] if y == 2025 else SCENARIOS["Base"][2025]
            curr_r = SCENARIOS["Base"][y]
            delta_bp = (curr_r - prev_r) * 100
            sector_d = (delta_bp / 25) * NII_SENS
            bank_d   = sector_d * (k24_nii[bank] / sector_nii)
            vol_d    = k24_nii[bank] * 0.02
            nii_curr = nii_prev + bank_d + vol_d
            fc_vals.append(nii_curr)
            nii_prev = nii_curr

        bridge_x = [2024, 2025, 2026]
        ax_nii.plot(bridge_x, fc_vals, color=bank_color_mpl(bank),
                    linewidth=1.5, linestyle="--", marker="D", markersize=4)
        ax_nii.annotate(bank.replace(" Bank",""),
                        xy=(2026, fc_vals[-1]),
                        xytext=(3, 0), textcoords="offset points",
                        fontsize=7, color=bank_color_mpl(bank))

    ax_nii.axvline(2024.5, linestyle=":", color="#475569", linewidth=1)
    ax_nii.set_title("NII Forecast — Base Scenario (€m)", fontsize=9.5)
    ax_nii.set_xticks([2022,2023,2024,2025,2026])
    ax_nii.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"€{v:,.0f}m"))
    ax_nii.grid(alpha=0.3)

    fig.tight_layout(pad=1.5)
    stream9 = fig_to_stream(fig, dpi=150)

add_image(sl9, stream9, 0.3, 1.2, 12.7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — STRESS TEST
# ════════════════════════════════════════════════════════════════════════════
sl10 = add_slide()
slide_header(sl10, "EBA-Style Stress Test",
             "Adverse scenario: +200bp CoR · −15% loans · −50bp NIM  |  Piraeus breaches 10.5% SREP floor")
footer(sl10, "10 / 12")

SREP = 10.5

def stress_bank_pptx(bank):
    row = kpis[kpis.bank==bank][kpis.year==2024].iloc[0]
    nii = row["nii"]; assets = row["total_assets"]; loans = row["loans"]
    impairment = row["impairment"]; opincome = row["operating_income"]
    opex = row["operating_expenses"]; net_profit = row["net_profit"]; cet1 = row["cet1"]
    nim = row["nim"] / 100
    nim_adv = nim - 50/10000
    loans_adv = loans * 0.85
    nii_adv = nim_adv * assets * 0.85
    other_income = opincome - nii
    ppop_adv = nii_adv + other_income + opex
    cor_adv = abs(impairment)/loans + 200/10000
    imp_adv = -(cor_adv * loans_adv)
    pbt_base = opincome + opex + impairment
    pbt_adv  = ppop_adv + imp_adv
    tr = max(0.15, min(0.35, 1 - net_profit/pbt_base)) if pbt_base else 0.24
    stressed_profit = pbt_adv * (1-tr)
    rwa = loans * 0.50
    stressed_cet1 = cet1 + (stressed_profit - net_profit) / rwa * 100
    return {"base_profit": net_profit, "adv_profit": stressed_profit,
            "base_cet1": cet1, "adv_cet1": stressed_cet1}

with mpl_context():
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4.5))
    fig.patch.set_facecolor("#0a1628")

    results = {b: stress_bank_pptx(b) for b in BANKS}

    # Left: Net profit baseline vs adverse
    x = np.arange(len(BANKS))
    w = 0.35
    base_p = [results[b]["base_profit"] for b in BANKS]
    adv_p  = [results[b]["adv_profit"]  for b in BANKS]
    ax1.bar(x-w/2, base_p, w, color=[bank_color_mpl(b) for b in BANKS], alpha=0.85, label="Baseline 2024")
    ax1.bar(x+w/2, adv_p,  w, color=[bank_color_mpl(b) for b in BANKS], alpha=0.4,  label="Stressed")
    ax1.set_title("Net Profit: Baseline vs Stressed (€m)", fontsize=9.5)
    ax1.set_xticks(x)
    ax1.set_xticklabels(BANKS, fontsize=8)
    ax1.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"€{v:,.0f}m"))
    ax1.legend(fontsize=8, framealpha=0)
    ax1.grid(axis="y", alpha=0.3)

    # Right: CET1 baseline vs adverse with SREP line
    base_c = [results[b]["base_cet1"] for b in BANKS]
    adv_c  = [results[b]["adv_cet1"]  for b in BANKS]
    ax2.bar(x-w/2, base_c, w, color=[bank_color_mpl(b) for b in BANKS], alpha=0.85, label="Baseline CET1")
    bar_colors = ["#ef4444" if results[b]["adv_cet1"] < SREP else bank_color_mpl(b) for b in BANKS]
    ax2.bar(x+w/2, adv_c, w, color=bar_colors, alpha=0.6, label="Stressed CET1")
    ax2.axhline(SREP, color="#f59e0b", linestyle="--", linewidth=2, label=f"SREP floor {SREP}%")
    ax2.set_title("CET1: Baseline vs Stressed (%)", fontsize=9.5)
    ax2.set_xticks(x)
    ax2.set_xticklabels(BANKS, fontsize=8)
    ax2.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"{v:.1f}%"))
    ax2.legend(fontsize=8, framealpha=0)
    ax2.grid(axis="y", alpha=0.3)
    # Annotate breach
    for i, b in enumerate(BANKS):
        adv = results[b]["adv_cet1"]
        color = "#f87171" if adv < SREP else "#86efac"
        ax2.text(i+w/2, adv+0.15, f"{adv:.1f}%", ha="center", fontsize=7.5, color=color)

    fig.tight_layout(pad=1.5)
    stream10 = fig_to_stream(fig, dpi=150)

add_image(sl10, stream10, 0.3, 1.25, 12.7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — VALUATION FRAMEWORK
# ════════════════════════════════════════════════════════════════════════════
sl11 = add_slide()
slide_header(sl11, "Valuation Framework — P/B vs Justified P/B",
             "Gordon Growth model: Justified P/B = ROE / CoE  |  CoE 11% (Rf 3.5% + ERP 5.5% + CRP 2%)")
footer(sl11, "11 / 12")

val_data = [
    ("Eurobank",     "~1.3x", 16.9, 11.0, "~1.5x", "~15% discount",    "Quality premium warranted",      BCOLORS_HEX["Eurobank"]),
    ("NBG",          "~1.3x", 13.7, 11.0, "~1.2x", "≈ fair value",     "Total return via buybacks",       BCOLORS_HEX["NBG"]),
    ("Piraeus Bank", "~0.9x", 12.9, 11.0, "~1.2x", "~25% discount",    "CoR risk closes gap",             BCOLORS_HEX["Piraeus Bank"]),
    ("Alpha Bank",   "~0.8x",  8.2, 11.0, "~0.7x", "Slight premium",   "Re-rate needs ROE > CoE",         BCOLORS_HEX["Alpha Bank"]),
]

add_rect(sl11, 0.3, 1.2, 12.5, 0.35, DGREY)
hdrs = ["Bank","Market P/B","ROE 2024","CoE","Justified P/B","vs Market","Investment Signal"]
col_xs = [0.35, 1.85, 3.35, 4.7, 5.85, 7.2, 8.9]
col_ws = [1.45, 1.45, 1.3, 1.1, 1.3, 1.65, 3.5]
for j, (h, cx, cw) in enumerate(zip(hdrs, col_xs, col_ws)):
    add_text(sl11, h, cx, 1.22, cw-0.05, 0.3,
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, (bank, mkt, roe, coe, just, vs_mkt, signal, bclr) in enumerate(val_data):
    row_y = 1.65 + i*0.85
    row_bg = rgb(15,30,50) if i % 2 == 0 else DGREY
    add_rect(sl11, 0.3, row_y, 12.5, 0.8, row_bg)
    add_rect(sl11, 0.3, row_y, 0.06, 0.8, bclr)
    vals_row = [bank, mkt, f"{roe:.1f}%", f"{coe:.1f}%", just, vs_mkt, signal]
    for j, (val, cx, cw) in enumerate(zip(vals_row, col_xs, col_ws)):
        align = PP_ALIGN.LEFT if j in (0, 6) else PP_ALIGN.CENTER
        is_bank = j == 0
        add_text(sl11, val, cx+0.08, row_y+0.18, cw-0.12, 0.45,
                 size=10 if not is_bank else 11,
                 bold=is_bank,
                 color=bclr if is_bank else WHITE,
                 align=align)

add_text(sl11,
    "Methodology: CoE = Risk-free rate (3.5%) + Beta (1.0) × Equity Risk Premium (5.5%) + Country Risk Premium (2.0%) = 11.0%  |  "
    "Market P/B approximate end-2024 prices. Justified P/B = ROE / CoE (Gordon Growth, g=0).",
    0.3, 5.15, 12.5, 0.5, size=8, color=LGREY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — METHODOLOGY
# ════════════════════════════════════════════════════════════════════════════
sl12 = add_slide()
slide_header(sl12, "Methodology & Data Quality",
             "All figures sourced from 12 audited annual report PDFs  |  No paid data vendors")
footer(sl12, "12 / 12")

method_cols = [
    ("Data Extraction",
     ["12 annual report PDFs",
      "pdfplumber (Python) extraction",
      "Manual cross-validation per figure",
      "Source page cited for every KPI"]),
    ("Data Quality",
     ["26 pytest assertions",
      "KPIs re-computed from CSVs",
      "BS identity check (Assets = L + E)",
      "NPE ratios: 12/12 PDF-verified"]),
    ("Financial Models",
     ["5-step banking DuPont",
      "CAMELS scoring (1–5 scale)",
      "Z-score peer benchmarking",
      "EBA-style single-period stress"]),
    ("Deliverables",
     ["SQLite DB + 3 processed CSVs",
      "Plotly/sql.js browser dashboard",
      "Streamlit multi-page app (live)",
      "Excel model + this deck"]),
]

for ci, (title, bullets) in enumerate(method_cols):
    cx = 0.35 + ci * 3.2
    add_rect(sl12, cx, 1.3, 3.0, 0.35, DGREY)
    add_text(sl12, title, cx+0.1, 1.32, 2.8, 0.3, size=10, bold=True, color=WHITE)
    for bi, b in enumerate(bullets):
        add_text(sl12, f"• {b}", cx+0.1, 1.75 + bi*0.45, 2.85, 0.42,
                 size=9, color=LGREY)

add_rect(sl12, 0.35, 5.0, 12.5, 0.04, BCOLORS_HEX["Eurobank"])
add_text(sl12,
    "Known limitations: Annual data only (no quarterly)  ·  NIM uses year-end assets (proxy)  ·  "
    "Market P/B values are end-2024 estimates  ·  RWA proxy = loans × 50%  ·  Not investment advice.",
    0.35, 5.1, 12.5, 0.45, size=8.5, color=LGREY)

add_text(sl12,
    "GitHub:  github.com/papastergiousp-maker/greek-banking-sector-analysis",
    0.35, 5.65, 8, 0.35, size=9, color=BCOLORS_HEX["Eurobank"])
add_text(sl12,
    "Live app:  greek-banking-sector-analysis.streamlit.app",
    0.35, 6.05, 8, 0.35, size=9, color=BCOLORS_HEX["Eurobank"])


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Written: {OUT}  ({len(prs.slides)} slides)")
